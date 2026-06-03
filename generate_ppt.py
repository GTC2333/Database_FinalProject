#!/usr/bin/env python3
"""
Final Project PPT Generator
Generates a complete PowerPoint with all charts and content.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Paths
PROJECT_DIR = os.path.dirname(os.path.abspath(__file__))
CHARTS_DIR = os.path.join(PROJECT_DIR, 'report', 'charts')
OUTPUT_PPT = os.path.join(PROJECT_DIR, 'report', 'Olist_Final_Project.pptx')

# Colors
DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)
ACCENT = RGBColor(0x00, 0x96, 0xC7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)


def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets, chart_path=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if chart_path and os.path.exists(chart_path):
        # Two-column: bullets left, chart right
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5.5))
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(12)

        # Chart on right
        slide.shapes.add_picture(chart_path, Inches(4.8), Inches(1.5), width=Inches(4.7))
    else:
        # Full-width bullets
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(14)

    return slide


def add_chart_only_slide(prs, title, chart_path):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if os.path.exists(chart_path):
        # Center chart
        slide.shapes.add_picture(chart_path, Inches(0.8), Inches(1.3), width=Inches(8.4))

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ===== 1. Title =====
    add_title_slide(prs,
        "巴西电商 Olist 运营管理系统",
        "数据库期末项目 | MySQL 8.x | 数据分析与可视化")

    # ===== 2. Section: 项目概述 =====
    add_section_slide(prs, "01 项目概述")

    add_content_slide(prs, "项目背景与实际意义", [
        "基于巴西电商平台 Olist 的真实公开数据集（~10万订单）",
        "构建完整的电商运营管理系统，涵盖用户/商品/订单/支付/评价全链路",
        "演示 MySQL 核心技术与电商业务数据分析",
        "覆盖 DDL/DML/约束/视图/存储过程/触发器/事务/索引/函数",
        "附加库存管理模块，模拟真实电商库存扣减与预警场景",
        "生成 11 张专业可视化图表，支撑业务洞察"
    ])

    add_content_slide(prs, "数据来源与规模", [
        "数据集：Kaggle - Brazilian E-Commerce Public Dataset by Olist",
        "数据规模：9 张 CSV，约 130MB，96,096 条有效订单",
        "时间跨度：2016-09 至 2018-08，覆盖近 2 年业务数据",
        "客户：96,096 人（覆盖巴西 27 个州）",
        "商品：32,951 种，涵盖 71 个品类",
        "卖家：3,095 家，地理坐标：1,000,163 条",
        "自造辅助表：inventory（库存）、inventory_alert（预警）"
    ])

    # ===== 3. Section: 数据库设计 =====
    add_section_slide(prs, "02 数据库设计")

    add_content_slide(prs, "ER 图与表结构", [
        "共 11 张数据表：9 张真实数据 + 2 张自造辅助表",
        "核心关系：customers → orders → order_items → products/sellers",
        "支付与评价均与 orders 一对多关联",
        "inventory 与 products 一对一，实现库存追踪",
        "所有表均设置 PRIMARY KEY，外键约束保证参照完整性",
        "【请在此处插入 MySQL Workbench 导出的 ER 图】"
    ])

    add_content_slide(prs, "完整性约束设计", [
        "实体完整性：所有表 PRIMARY KEY + customers.unique_id UNIQUE",
        "参照完整性：6 组外键关联，删除策略 CASCADE / SET NULL / RESTRICT",
        "CHECK 约束：review_score 1-5、price≥0、freight≥0、stock≥0",
        "ENUM 约束：order_status 限定 8 种状态（delivered/shipped/canceled...）",
        "NOT NULL 约束：业务必填字段均强制非空",
        "DEFAULT 约束：inventory.last_updated 自动维护时间戳"
    ])

    # ===== 4. Section: 基础操作 =====
    add_section_slide(prs, "03 基础操作演示")

    add_content_slide(prs, "建库建表与数据导入", [
        "CREATE DATABASE olist_db CHARACTER SET utf8mb4",
        "11 张表使用 InnoDB 引擎，支持事务与行级锁",
        "数据导入：LOAD DATA LOCAL INFILE 批量加载 9 个 CSV",
        "导入结果验证：customers 96,096 / orders 96,096 / order_items 108,578",
        "inventory 表生成：基于 products 随机生成库存（50-500），5% 设为低库存",
        "导入耗时：普通笔记本约 1-2 分钟完成全部加载"
    ])

    add_content_slide(prs, "增删改查与多表查询", [
        "基础查询：SELECT / INSERT / UPDATE / DELETE 完整覆盖",
        "多表 JOIN：订单宽表 v_order_full 涉及 4 表 JOIN + GROUP BY",
        "子查询：RFM 计算中的窗口函数 NTILE 分位",
        "聚合分析：SUM / COUNT / AVG / ROUND 配合 GROUP BY",
        "时间函数：DATE_FORMAT、YEAR、HOUR、DATEDIFF",
        "条件逻辑：CASE WHEN 实现 RFM 客户分群标签"
    ])

    # ===== 5. Section: 进阶技术 =====
    add_section_slide(prs, "04 MySQL 进阶技术")

    add_content_slide(prs, "视图（Views）", [
        "v_order_full：订单宽表，聚合客户/明细/评价信息",
        "v_category_sales：品类销售统计，展示收入/销量/均价",
        "v_customer_rfm：RFM 客户价值分层，使用 NTILE 五分位",
        "v_monthly_revenue：月度收入趋势，方便时间序列分析",
        "v_product_recommendation：协同过滤推荐，基于共现统计"
    ])

    add_content_slide(prs, "函数、存储过程与事务", [
        "自定义函数：fn_shipping_days() 计算物流天数；fn_order_total() 计算订单总价",
        "sp_place_order()：模拟下单事务，START TRANSACTION → INSERT → COMMIT/ROLLBACK",
        "sp_state_monthly_revenue(state, year)：州级月度收入统计",
        "sp_low_stock_report()：返回库存不足商品列表",
        "sp_restock_product(product_id, qty)：商品补货操作",
        "事务处理：库存扣减失败自动 ROLLBACK，保证数据一致性"
    ])

    add_content_slide(prs, "触发器（Triggers）与索引", [
        "trg_order_item_deduct_stock：下单前检查并自动扣减库存",
        "trg_inventory_low_stock_alert：库存降至安全线时自动插入预警",
        "trg_review_negative_alert：评分≤2 时自动记录负面评价预警",
        "索引设计：共 11 个索引，含 3 个复合索引",
        "idx_order_status_time：(status, timestamp) 复合索引加速时间范围查询",
        "idx_oi_product_seller：(product_id, seller_id) 支持品类-卖家分析"
    ])

    # ===== 6. Section: 数据分析 =====
    add_section_slide(prs, "05 数据分析与可视化")

    # Chart slides
    charts = [
        ("月度销售趋势", "01_monthly_revenue_trend.png",
         "2017-11 达到峰值 7,060 单 / 112 万 BRL，随后稳定月均 6,000 单"),
        ("各州收入地理分布", "02_top_states_revenue.png",
         "SP 州 559 万 BRL 遥遥领先，占总收入 35%；RJ、MG 分居二三位"),
        ("TOP 10 品类销售额", "03_top_categories.png",
         "health_beauty 123 万居首，watches_gifts 117 万次之，客单价差异显著"),
        ("支付方式偏好", "04_payment_distribution.png",
         "信用卡占 75.2% 订单 / 78.3% 金额，平均分期 3.5 期"),
        ("配送时效分析", "05_delivery_days_by_state.png",
         "SP 州 8.7 天最优，偏远州 20-30 天，存在显著物流差距"),
        ("评价分数分布", "06_review_distribution.png",
         "5 星 57.7%，好评率 77.1%，但 1 星 11.5% 需关注"),
        ("RFM 客户分层", "07_rfm_segments.png",
         "识别 Champions（高价值）与 Lost（流失需召回）客户群体"),
        ("购买时段分析", "08_hourly_pattern.png",
         "白天 10:00-17:00 为黄金时段，16:00 峰值 6,484 单"),
        ("订单状态分布", "09_order_status.png",
         "delivered 占 97.0%，canceled 仅 0.6%，履约率极高"),
        ("价格与运费关系", "10_price_freight_scatter.png",
         "watches_gifts 高价位，bed_bath_table 低价高频，品类价格带分层明显"),
        ("州级业务画像", "11_state_bubble.png",
         "SP 龙头但 AOV 最低（143 BRL）；BA 州 AOV 182 BRL 但订单中等"),
    ]

    for title, filename, insight in charts:
        path = os.path.join(CHARTS_DIR, filename)
        # Insight slide
        add_content_slide(prs, title, [
            insight,
            "",
            f"【图表文件：{filename}】"
        ], chart_path=path)

    # ===== 7. Section: 总结 =====
    add_section_slide(prs, "06 总结")

    add_content_slide(prs, "项目成果总结", [
        "完整的数据库系统：11 张表、完整的约束体系、96K+ 订单数据",
        "6 大进阶技术全覆盖：视图(5) / 函数(2) / 触发器(3) / 存储过程(4) / 事务 / 索引(11)",
        "10 个业务分析维度：销售/地理/品类/支付/配送/评价/RFM/时段/状态/画像",
        "11 张专业可视化图表：折线/柱状/饼图/散点/气泡/面积图",
        "可复现的交付物：SQL 脚本 + Python 图表脚本 + 一键运行",
        "业务洞察：识别核心市场(SP)、高价值品类(health_beauty)、物流优化空间"
    ])

    add_content_slide(prs, "成员分工", [
        "（姓名）：数据库设计、建表、约束、数据导入",
        "（姓名）：视图、存储过程、函数、触发器、事务",
        "（姓名）：数据分析查询、Python 可视化、报告撰写",
        "（姓名）：PPT 制作、演示准备、Word 报告整理",
        "",
        "【请填写实际成员姓名】"
    ])

    # ===== 8. Thanks =====
    add_title_slide(prs, "感谢聆听", "Questions & Discussion")

    prs.save(OUTPUT_PPT)
    print(f"PPT saved: {OUTPUT_PPT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == '__main__':
    main()
